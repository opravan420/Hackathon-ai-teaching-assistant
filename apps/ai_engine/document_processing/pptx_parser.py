import pptx
from .base import BaseDocumentParser
from .exceptions import DocumentExtractionError, NoExtractableTextError, InvalidDocumentError

class PPTXParser(BaseDocumentParser):
    """Parser for extracting text from PPTX presentation files using python-pptx."""

    def extract_text(self, file_path: str) -> str:
        try:
            prs = pptx.Presentation(file_path)
        except Exception as e:
            raise InvalidDocumentError(f"Corrupted or invalid PPTX file: {str(e)}")

        try:
            extracted_slides = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text = [f"[Slide {i}]"]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            val = paragraph.text.strip()
                            if val:
                                slide_text.append(val)
                    
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_text:
                                slide_text.append(" | ".join(row_text))

                # Only include slides that have actual text content other than the header label
                if len(slide_text) > 1:
                    extracted_slides.append("\n".join(slide_text))

            full_text = "\n\n".join(extracted_slides).strip()
            if not full_text:
                raise NoExtractableTextError("PPTX contains no extractable text.")

            return full_text
        except NoExtractableTextError:
            raise
        except Exception as e:
            raise DocumentExtractionError(f"Error extracting text from PPTX: {str(e)}")
