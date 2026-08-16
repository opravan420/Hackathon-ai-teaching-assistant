import os
import shutil
import fitz
import docx
import pptx
from django.test import TestCase, Client
from django.urls import reverse
from django.contrib.auth import get_user_model
from django.core.files.uploadedfile import SimpleUploadedFile
from django.test.signals import template_rendered
from django.test.client import store_rendered_templates
from unittest.mock import MagicMock
from apps.accounts.models import TeacherProfile
from apps.ai_engine.models import Document
from apps.ai_engine.document_processing.exceptions import (
    UnsupportedFileTypeError,
    FileSizeLimitError,
    InvalidDocumentError,
    NoExtractableTextError
)
from apps.ai_engine.document_processing.service import DocumentService, normalize_text

# Workaround for Django 5.0 template context copying on Python 3.14
template_rendered.disconnect(store_rendered_templates)

User = get_user_model()

class DocumentProcessingTestCase(TestCase):
    def setUp(self):
        # Workaround for Django 5.0 template context copying on Python 3.14
        import django.test.client
        def custom_store(store, signal, sender, template, context, **kwargs):
            store.setdefault("templates", []).append(template)
        django.test.client.store_rendered_templates = custom_store

        # Create two teachers to test ownership bounds
        self.teacher_1 = User.objects.create_user(
            username='teacher_doc1', password='password123', role=User.TEACHER
        )
        self.profile_1 = TeacherProfile.objects.create(
            user=self.teacher_1, employee_id='EMP_D1', department='Physics'
        )

        self.teacher_2 = User.objects.create_user(
            username='teacher_doc2', password='password123', role=User.TEACHER
        )
        self.profile_2 = TeacherProfile.objects.create(
            user=self.teacher_2, employee_id='EMP_D2', department='Chem'
        )

        self.client1 = Client()
        self.client1.login(username='teacher_doc1', password='password123')

        self.client2 = Client()
        self.client2.login(username='teacher_doc2', password='password123')

        self.service = DocumentService()
        self.temp_dir = os.path.join(os.path.dirname(__file__), 'temp_test_fixtures')
        os.makedirs(self.temp_dir, exist_ok=True)

    def tearDown(self):
        if os.path.exists(self.temp_dir):
            shutil.rmtree(self.temp_dir)
        # Clean up any document files created in media
        for doc in Document.objects.all():
            if doc.stored_file and os.path.exists(doc.stored_file.path):
                try:
                    os.remove(doc.stored_file.path)
                except Exception:
                    pass

    # 1. Helper to construct temporary files
    def make_file_path(self, filename):
        return os.path.join(self.temp_dir, filename)

    # 2. Text normalization test
    def test_text_normalization(self):
        dirty = "  Line 1 \r\n\r\n\r\n Line 2 \r\n   \r\n  Line 3  "
        expected = "Line 1\n\nLine 2\n\nLine 3"
        self.assertEqual(normalize_text(dirty), expected)

    # 3. Plain text and Hindi/English UTF-8 checks
    def test_txt_extraction_hindi_english(self):
        path = self.make_file_path("test.txt")
        text_content = "Deadlock is a situation...\nडेडलॉक एक ऐसी स्थिति है..."
        with open(path, "w", encoding="utf-8") as f:
            f.write(text_content)

        with open(path, "rb") as f:
            upload_file = SimpleUploadedFile("test.txt", f.read())

        doc = self.service.process_document(self.teacher_1, upload_file)
        self.assertEqual(doc.extraction_status, Document.SUCCESS)
        self.assertEqual(doc.file_type, "TXT")
        self.assertIn("डेडलॉक एक ऐसी स्थिति है", doc.extracted_text)
        self.assertIn("Deadlock is a situation", doc.extracted_text)
        self.assertEqual(doc.character_count, len(doc.extracted_text))

    # 4. Empty plain text check
    def test_empty_txt_raises_no_text_error(self):
        upload_file = SimpleUploadedFile("empty.txt", b"")
        with self.assertRaises(NoExtractableTextError):
            self.service.process_document(self.teacher_1, upload_file)
        
        # Verify db status logged as NO_EXTRACTABLE_TEXT
        doc = Document.objects.get(original_filename="empty.txt")
        self.assertEqual(doc.extraction_status, Document.NO_EXTRACTABLE_TEXT)

    # 5. PDF extraction check
    def test_pdf_extraction_success(self):
        path = self.make_file_path("test.pdf")
        doc_pdf = fitz.open()
        page = doc_pdf.new_page()
        page.insert_text((50, 50), "Hello world from PDF")
        doc_pdf.save(path)
        doc_pdf.close()

        with open(path, "rb") as f:
            upload_file = SimpleUploadedFile("test.pdf", f.read())

        doc = self.service.process_document(self.teacher_1, upload_file)
        self.assertEqual(doc.extraction_status, Document.SUCCESS)
        self.assertEqual(doc.file_type, "PDF")
        self.assertIn("Hello world from PDF", doc.extracted_text)

    # 6. Invalid PDF file handling
    def test_invalid_pdf_raises_error(self):
        upload_file = SimpleUploadedFile("corrupted.pdf", b"random corrupted non-pdf bytes")
        with self.assertRaises(InvalidDocumentError):
            self.service.process_document(self.teacher_1, upload_file)
        
        doc = Document.objects.get(original_filename="corrupted.pdf")
        self.assertEqual(doc.extraction_status, Document.FAILED)

    # 7. DOCX extraction check
    def test_docx_extraction_success(self):
        path = self.make_file_path("test.docx")
        doc_docx = docx.Document()
        doc_docx.add_paragraph("Hello world from DOCX")
        doc_docx.save(path)

        with open(path, "rb") as f:
            upload_file = SimpleUploadedFile("test.docx", f.read())

        doc = self.service.process_document(self.teacher_1, upload_file)
        self.assertEqual(doc.extraction_status, Document.SUCCESS)
        self.assertEqual(doc.file_type, "DOCX")
        self.assertEqual(doc.extracted_text, "Hello world from DOCX")

    # 8. Invalid DOCX file handling
    def test_invalid_docx_raises_error(self):
        upload_file = SimpleUploadedFile("corrupted.docx", b"random corrupted docx bytes")
        with self.assertRaises(InvalidDocumentError):
            self.service.process_document(self.teacher_1, upload_file)
        
        doc = Document.objects.get(original_filename="corrupted.docx")
        self.assertEqual(doc.extraction_status, Document.FAILED)

    # 9. PPTX extraction check
    def test_pptx_extraction_success(self):
        path = self.make_file_path("test.pptx")
        prs = pptx.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(0, 0, 100, 100)
        tf = txBox.text_frame
        tf.text = "Hello world from PPTX slide"
        prs.save(path)

        with open(path, "rb") as f:
            upload_file = SimpleUploadedFile("test.pptx", f.read())

        doc = self.service.process_document(self.teacher_1, upload_file)
        self.assertEqual(doc.extraction_status, Document.SUCCESS)
        self.assertEqual(doc.file_type, "PPTX")
        self.assertIn("Hello world from PPTX slide", doc.extracted_text)
        self.assertIn("[Slide 1]", doc.extracted_text)

    # 10. Invalid PPTX file handling
    def test_invalid_pptx_raises_error(self):
        upload_file = SimpleUploadedFile("corrupted.pptx", b"random corrupted pptx bytes")
        with self.assertRaises(InvalidDocumentError):
            self.service.process_document(self.teacher_1, upload_file)
        
        doc = Document.objects.get(original_filename="corrupted.pptx")
        self.assertEqual(doc.extraction_status, Document.FAILED)

    # 11. Unsupported format (Extension check)
    def test_unsupported_extension_raises_error(self):
        upload_file = SimpleUploadedFile("image.png", b"png bytes")
        with self.assertRaises(UnsupportedFileTypeError):
            self.service.process_document(self.teacher_1, upload_file)

    # 12. Legacy .ppt format explicit rejection
    def test_legacy_ppt_explicit_rejection(self):
        upload_file = SimpleUploadedFile("lecture.ppt", b"legacy ppt bytes")
        with self.assertRaises(UnsupportedFileTypeError) as context:
            self.service.process_document(self.teacher_1, upload_file)
        self.assertIn("legacy .ppt format is not supported", str(context.exception))

    # 13. File too large limit check
    def test_file_too_large_rejection(self):
        # Create a mock file with size exceeding 25MB
        dummy_file = MagicMock()
        dummy_file.name = "huge.pdf"
        dummy_file.size = 26 * 1024 * 1024 # 26 MB
        
        with self.assertRaises(FileSizeLimitError):
            self.service.process_document(self.teacher_1, dummy_file)

    # 14. Teacher ownership bounds check
    def test_teacher_ownership(self):
        # Create document for teacher 1
        path = self.make_file_path("teacher1_doc.txt")
        with open(path, "w", encoding="utf-8") as f:
            f.write("Teacher 1 secret content")
        with open(path, "rb") as f:
            upload_file = SimpleUploadedFile("teacher1_doc.txt", f.read())
        
        doc = self.service.process_document(self.teacher_1, upload_file)
        
        # Verify teacher 1 query shows it
        t1_docs = Document.objects.filter(teacher=self.teacher_1)
        self.assertIn(doc, t1_docs)
        
        # Verify teacher 2 query does NOT show it
        t2_docs = Document.objects.filter(teacher=self.teacher_2)
        self.assertNotIn(doc, t2_docs)
